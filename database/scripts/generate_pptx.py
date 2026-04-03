"""Generate Marketing Presentation for BanK Data Quality Monitor — FRENCH VERSION"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# === COLOR PALETTE ===
DARK_BG = RGBColor(0x0B, 0x1D, 0x3A)
ACCENT_BLUE = RGBColor(0x00, 0x7B, 0xFF)
ACCENT_TEAL = RGBColor(0x00, 0xC9, 0xA7)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xB0, 0xBC, 0xCE)
GOLD = RGBColor(0xFF, 0xC1, 0x07)
ORANGE = RGBColor(0xFF, 0x6B, 0x35)
CARD_BG = RGBColor(0x11, 0x2B, 0x4A)
SECTION_BG = RGBColor(0x08, 0x15, 0x2B)


def add_dark_bg(slide, color=DARK_BG):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    sp = bg._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_accent_bar(slide, left, top, width, height, color=ACCENT_BLUE):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_card(slide, left, top, width, height, color=CARD_BG):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    return card


def add_bullet_card(slide, left, top, width, height, title, bullets, title_color=ACCENT_TEAL):
    add_card(slide, left, top, width, height)
    add_text_box(slide, left + Inches(0.3), top + Inches(0.15), width - Inches(0.6), Inches(0.5),
                 title, font_size=16, color=title_color, bold=True)
    txBox = slide.shapes.add_textbox(left + Inches(0.3), top + Inches(0.6), width - Inches(0.6), height - Inches(0.8))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Segoe UI"
        p.space_after = Pt(4)


# ============================================================
# SLIDE 1 — PAGE DE TITRE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide, SECTION_BG)

add_accent_bar(slide, Inches(1.5), Inches(2.3), Inches(1.0), Inches(0.06), ACCENT_TEAL)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10), Inches(1.2),
             "BanK Data Quality Monitor", font_size=44, color=WHITE, bold=True)

add_text_box(slide, Inches(1.5), Inches(3.6), Inches(8), Inches(0.6),
             "Plateforme Intelligente de Qualite des Donnees Bancaires & Conformite FATCA",
             font_size=20, color=ACCENT_TEAL)

add_text_box(slide, Inches(1.5), Inches(4.4), Inches(8), Inches(1.0),
             "Solution complete pour detecter, suivre et corriger les anomalies\ndans les systemes d'information client tout en assurant la conformite reglementaire.",
             font_size=14, color=LIGHT_GRAY)

add_accent_bar(slide, 0, Inches(6.4), prs.slide_width, Inches(1.1), CARD_BG)
stats = [("18+", "Pages"), ("90+", "Points API"), ("20+", "Tables BD"), ("4", "Roles"), ("3", "Types Client")]
for i, (num, lbl) in enumerate(stats):
    x = Inches(1.5 + i * 2.2)
    add_text_box(slide, x, Inches(6.5), Inches(1.5), Inches(0.5), num, font_size=28, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, Inches(6.95), Inches(1.5), Inches(0.3), lbl, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 2 — LE DEFI
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), ACCENT_TEAL)
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(6), Inches(0.6),
             "Le Defi", font_size=32, color=WHITE, bold=True)
add_text_box(slide, Inches(1.1), Inches(1.1), Inches(10), Inches(0.5),
             "Les banques font face a des defis critiques en matiere de qualite des donnees et de conformite",
             font_size=15, color=LIGHT_GRAY)

challenges = [
    ("Lacunes de Qualite", "Enregistrements clients incoherents, incomplets ou en double entre les agences, source d'erreurs operationnelles et de risques reglementaires.", ORANGE),
    ("Pression FATCA", "Le filtrage manuel des indices americains est lent, sujet aux erreurs et ne peut pas passer a l'echelle sur des centaines d'agences.", ORANGE),
    ("Manque de Gouvernance", "Les corrections appliquees sans piste d'audit creent des vulnerabilites et des non-conformites.", ORANGE),
    ("Systemes Cloisonnes", "Les donnees du core banking deconnectees des outils de suivi qualite — les problemes restent non detectes pendant des mois.", ORANGE),
    ("Aucun Progres Mesurable", "Sans KPI ni tableaux de bord, la direction ne peut pas suivre ni prouver l'amelioration de la qualite des donnees.", ORANGE),
    ("Processus Manuels", "Les corrections necessitent une re-saisie manuelle entre systemes, gaspillant des milliers d'heures par an.", ORANGE),
]

for i, (title, desc, color) in enumerate(challenges):
    col = i % 3
    row = i // 3
    x = Inches(0.8 + col * 4.0)
    y = Inches(1.8 + row * 2.6)
    add_card(slide, x, y, Inches(3.7), Inches(2.2))
    add_accent_bar(slide, x, y, Inches(3.7), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.1), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.8), Inches(3.1), Inches(1.2),
                 desc, font_size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE 3 — VUE D'ENSEMBLE DE LA SOLUTION
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), ACCENT_TEAL)
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.6),
             "Notre Solution — Plateforme Qualite de Bout en Bout", font_size=30, color=WHITE, bold=True)

steps = [
    ("DETECTER", "Detection automatique\ndes anomalies via\nregles configurables", ACCENT_BLUE),
    ("TICKET", "Creation automatique\nde tickets et\nassignation", RGBColor(0x6C, 0x5C, 0xE7)),
    ("CORRIGER", "Correction guidee\navec automatisation\nRPA", ACCENT_TEAL),
    ("VALIDER", "Validation 4-yeux\npar le superviseur", GOLD),
    ("RECONCILIER", "Synchronisation\net verification\ncore banking", ORANGE),
    ("RAPPORTER", "Tableaux de bord\nKPI et rapports\nreglementaires", RGBColor(0xE8, 0x4D, 0x8A)),
]

for i, (title, desc, color) in enumerate(steps):
    x = Inches(0.6 + i * 2.1)
    y = Inches(1.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.8), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    if i < len(steps) - 1:
        add_text_box(slide, x + Inches(1.8), y + Inches(0.1), Inches(0.3), Inches(0.5),
                     ">", font_size=20, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    add_text_box(slide, x, y + Inches(0.9), Inches(1.8), Inches(1.0),
                 desc, font_size=11, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.04), ACCENT_TEAL)
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(5), Inches(0.5),
             "Differenciateurs Cles", font_size=20, color=ACCENT_TEAL, bold=True)

diffs = [
    ("Integration Core Banking Directe", "Connexion directe a Informix/BSIC via JDBC — aucun middleware requis"),
    ("Architecture Multi-Agences", "Vues par agence avec supervision du siege — scalable de 1 a 500+ agences"),
    ("Pipeline d'Automatisation Complet", "De la detection a la reconciliation — de bout en bout sans re-saisie manuelle"),
    ("Pret pour la Reglementation", "Filtrage FATCA + validation 4-yeux + piste d'audit complete"),
]

for i, (title, desc) in enumerate(diffs):
    col = i % 2
    row = i // 2
    x = Inches(0.8 + col * 6.0)
    y = Inches(4.6 + row * 1.2)
    add_accent_bar(slide, x, y + Inches(0.05), Inches(0.06), Inches(0.35), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.2), y, Inches(5.5), Inches(0.35),
                 title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.4), Inches(5.5), Inches(0.5),
                 desc, font_size=11, color=LIGHT_GRAY)


# ============================================================
# SLIDE 4 — CONTROLE QUALITE AUTOMATISE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), ACCENT_BLUE)
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.6),
             "Controle Qualite Automatise des Donnees", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(1.1), Inches(1.05), Inches(10), Inches(0.4),
             '"Detectez les problemes de qualite avant les regulateurs."', font_size=14, color=ACCENT_TEAL, bold=True)

features = [
    ("Detection Multi-Types",
     ["Anomalies clients particuliers", "Anomalies clients entreprises", "Anomalies institutionnels",
      "Couverture 100% du portefeuille"]),
    ("Moteur de Regles Configurable",
     ["Regles de validation via editeur SQL", "Requetes de detection pre-construites", "Activation/desactivation des regles",
      "Aucune dependance fournisseur"]),
    ("Suivi en Temps Reel",
     ["Scan instantane des anomalies", "Suivi des statuts (detecte/en cours/corrige)", "Historique champ par champ",
      "Tracabilite complete"]),
    ("Operations en Masse",
     ["Traitement massif des anomalies", "Filtrage par agence", "Recherche par identifiant client",
      "Export PDF pour reporting"]),
]

for i, (title, bullets) in enumerate(features):
    x = Inches(0.8 + i * 3.1)
    add_bullet_card(slide, x, Inches(1.7), Inches(2.8), Inches(3.2), title, bullets, title_color=ACCENT_BLUE)

add_card(slide, Inches(0.8), Inches(5.2), Inches(11.7), Inches(1.6))
add_text_box(slide, Inches(1.2), Inches(5.4), Inches(4), Inches(0.4),
             "Couverture de Detection", font_size=16, color=WHITE, bold=True)

coverage = [
    ("Documents d'Identite Manquants", "Particuliers sans identification valide"),
    ("Enregistrement Incomplet", "Entreprises sans donnees d'immatriculation"),
    ("Validation d'Adresse", "Clients avec adresses incompletes ou invalides"),
    ("Controle Tel/Email", "Coordonnees manquantes ou mal formatees"),
]
for i, (title, desc) in enumerate(coverage):
    x = Inches(1.2 + i * 2.9)
    y = Inches(5.9)
    add_accent_bar(slide, x, y, Inches(0.06), Inches(0.6), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.2), y, Inches(2.5), Inches(0.3), title, font_size=12, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.3), Inches(2.5), Inches(0.4), desc, font_size=10, color=LIGHT_GRAY)


# ============================================================
# SLIDE 5 — CONFORMITE FATCA
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), GOLD)
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.6),
             "Conformite FATCA — Toujours Pret pour l'Audit", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(1.1), Inches(1.05), Inches(10), Inches(0.4),
             '"Conformite FATCA en pilote automatique — reduit l\'effort de filtrage manuel de 80%+"',
             font_size=14, color=GOLD, bold=True)

# Left — Filtrage des indices
add_card(slide, Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.0))
add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.4),
             "Filtrage Automatique des Indices US", font_size=18, color=GOLD, bold=True)

indicia = [
    ("Lieu de Naissance US", "Detection automatique des lieux de naissance en territoire americain"),
    ("Nationalite US", "Detection et verification du drapeau de nationalite americaine"),
    ("Adresse US", "Correspondance de motifs d'adresses americaines"),
    ("Numero de Telephone US", "Detection du format telephonique (+1, indicatifs US)"),
]
for i, (title, desc) in enumerate(indicia):
    y = Inches(2.5 + i * 0.95)
    add_accent_bar(slide, Inches(1.1), y + Inches(0.05), Inches(0.5), Inches(0.06), GOLD)
    add_text_box(slide, Inches(1.1), y + Inches(0.15), Inches(4.8), Inches(0.3), title, font_size=14, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.1), y + Inches(0.45), Inches(4.8), Inches(0.35), desc, font_size=11, color=LIGHT_GRAY)

# Right — Workflow et reporting
add_card(slide, Inches(6.6), Inches(1.7), Inches(5.9), Inches(2.3))
add_text_box(slide, Inches(6.9), Inches(1.9), Inches(5.3), Inches(0.4),
             "Workflow de Conformite a 4 Statuts", font_size=18, color=GOLD, bold=True)

statuses = [
    ("A verifier", "Signale pour revue", ORANGE),
    ("Confirme", "Personne US confirmee", RGBColor(0xE7, 0x4C, 0x3C)),
    ("Exclu", "Non US - valide", ACCENT_TEAL),
    ("En attente", "Info supplementaire", ACCENT_BLUE),
]
for i, (status, desc, color) in enumerate(statuses):
    x = Inches(6.9 + i * 1.4)
    y = Inches(2.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Inches(1.2), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(10)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER
    add_text_box(slide, x, y + Inches(0.55), Inches(1.2), Inches(0.3), desc, font_size=9, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Audit card
add_card(slide, Inches(6.6), Inches(4.2), Inches(5.9), Inches(2.5))
add_text_box(slide, Inches(6.9), Inches(4.4), Inches(5.3), Inches(0.4),
             "Piste d'Audit & Reporting", font_size=18, color=GOLD, bold=True)
audit_items = [
    "Journal d'audit FATCA complet avec historique des modifications",
    "Statistiques FATCA par agence pour les responsables",
    "Export PDF et transmission aux autorites fiscales",
    "Suivi des clients FATCA individuels et entreprises",
    "Classification et regroupement par indices",
    "Pret pour l'inspection IRS/reglementaire a tout moment",
]
txBox = slide.shapes.add_textbox(Inches(6.9), Inches(4.9), Inches(5.3), Inches(1.6))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(audit_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Segoe UI"
    p.space_after = Pt(3)


# ============================================================
# SLIDE 6 — GOUVERNANCE & WORKFLOW
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), ACCENT_TEAL)
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.6),
             "Validation 4-Yeux & Gouvernance", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(1.1), Inches(1.05), Inches(10), Inches(0.4),
             '"Gouvernance integree — aucune correction ne passe sans controle."', font_size=14, color=ACCENT_TEAL, bold=True)

wf_steps = [
    ("1. DETECTER", "Anomalie decouverte\nvia regle de validation"),
    ("2. TICKET", "Ticket de support\ncree automatiquement"),
    ("3. SOUMETTRE", "L'agent soumet\nune correction"),
    ("4. VALIDER", "Le superviseur\nrevise et approuve"),
    ("5. APPLIQUER", "Correction appliquee\nau systeme"),
    ("6. RECONCILIER", "Synchronise avec\nle core banking"),
]

for i, (title, desc) in enumerate(wf_steps):
    x = Inches(0.6 + i * 2.1)
    y = Inches(1.7)
    color = [ACCENT_BLUE, RGBColor(0x6C, 0x5C, 0xE7), ACCENT_TEAL, GOLD, ORANGE, RGBColor(0xE8, 0x4D, 0x8A)][i]

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.55), y, Inches(0.7), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"

    add_text_box(slide, x, y + Inches(0.8), Inches(1.8), Inches(0.35),
                 title.split(". ")[1], font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x, y + Inches(1.15), Inches(1.8), Inches(0.7),
                 desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

    if i < len(wf_steps) - 1:
        add_text_box(slide, x + Inches(1.65), y + Inches(0.15), Inches(0.5), Inches(0.5),
                     "-->", font_size=14, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# RBAC
add_accent_bar(slide, Inches(0.8), Inches(3.8), Inches(11.7), Inches(0.04), ACCENT_TEAL)
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(6), Inches(0.4),
             "Controle d'Acces Base sur les Roles (RBAC)", font_size=20, color=ACCENT_TEAL, bold=True)

roles = [
    ("ADMIN", "Acces systeme complet\nGestion utilisateurs & config\nCreation de regles\nToutes operations CRUD", ACCENT_BLUE),
    ("AUDITEUR", "Acces lecture & validation\nApprobation corrections (4-yeux)\nConsultation des rapports\nAutorite approbation FATCA", RGBColor(0x6C, 0x5C, 0xE7)),
    ("UTILISATEUR AGENCE", "Vues limitees a l'agence\nSoumission de corrections\nConsultation tickets agence\nDemande de validation", ACCENT_TEAL),
    ("UTILISATEUR", "Acces lecture de base\nConsultation propres donnees\nSoumission de demandes\nPermissions standard", LIGHT_GRAY),
]

for i, (role, perms, color) in enumerate(roles):
    x = Inches(0.8 + i * 3.1)
    y = Inches(4.5)
    add_card(slide, x, y, Inches(2.8), Inches(2.5))
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.15), y + Inches(0.2), Inches(2.5), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = role
    p.font.size = Pt(11)
    p.font.color.rgb = WHITE if color != LIGHT_GRAY else DARK_BG
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, x + Inches(0.2), y + Inches(0.8), Inches(2.4), Inches(1.5),
                 perms, font_size=10, color=LIGHT_GRAY)


# ============================================================
# SLIDE 7 — INTEGRATION CORE BANKING
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), ORANGE)
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.6),
             "Integration Core Banking & Reconciliation", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(1.1), Inches(1.05), Inches(10), Inches(0.4),
             '"Connexion directe a votre core banking — zero silo de donnees."', font_size=14, color=ORANGE, bold=True)

add_card(slide, Inches(0.8), Inches(1.7), Inches(5.8), Inches(5.0))
add_text_box(slide, Inches(1.1), Inches(1.9), Inches(5), Inches(0.4),
             "Architecture Systeme", font_size=18, color=ORANGE, bold=True)

systems = [
    ("Core Banking Informix (BSIC)", "Connexion JDBC directe — clients, comptes, adresses, contacts", ORANGE),
    ("Base de Donnees PostgreSQL", "Donnees qualite, anomalies, tickets, FATCA, journaux d'audit", ACCENT_BLUE),
    ("Keycloak SSO", "Authentification OAuth2 — connexion unique entreprise", ACCENT_TEAL),
    ("Flowable BPMN", "Moteur d'automatisation de workflows standard", RGBColor(0x6C, 0x5C, 0xE7)),
    ("UiPath RPA", "Automatisation robotique des processus de correction", GOLD),
    ("Cache Redis", "Optimisation des performances pour tableaux de bord temps reel", RGBColor(0xE8, 0x4D, 0x8A)),
]

for i, (sys_name, desc, color) in enumerate(systems):
    y = Inches(2.5 + i * 0.7)
    add_accent_bar(slide, Inches(1.1), y + Inches(0.05), Inches(0.06), Inches(0.4), color)
    add_text_box(slide, Inches(1.3), y, Inches(5), Inches(0.3), sys_name, font_size=12, color=color, bold=True)
    add_text_box(slide, Inches(1.3), y + Inches(0.28), Inches(5), Inches(0.3), desc, font_size=10, color=LIGHT_GRAY)

# Reconciliation
add_card(slide, Inches(6.9), Inches(1.7), Inches(5.6), Inches(2.3))
add_text_box(slide, Inches(7.2), Inches(1.9), Inches(5), Inches(0.4),
             "Moteur de Reconciliation", font_size=18, color=ORANGE, bold=True)
recon_items = [
    "Synchronisation bidirectionnelle (CBS <-> Application)",
    "Reconciliation unitaire et par lots (jusqu'a 50)",
    "Score de correspondance par enregistrement",
    "Suivi des echecs et mecanisme de reessai",
    "Historique et statistiques de reconciliation",
]
txBox = slide.shapes.add_textbox(Inches(7.2), Inches(2.4), Inches(5), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(recon_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Segoe UI"
    p.space_after = Pt(3)

# RPA
add_card(slide, Inches(6.9), Inches(4.2), Inches(5.6), Inches(2.5))
add_text_box(slide, Inches(7.2), Inches(4.4), Inches(5), Inches(0.4),
             "Automatisation RPA", font_size=18, color=GOLD, bold=True)
rpa_items = [
    "Integration UiPath pour corrections automatisees",
    "Suivi des jobs en temps reel (refresh auto 10s)",
    "Auto-reparation avec reessai automatique",
    "Tableau de bord statistiques pour suivi du ROI",
    "Callbacks de statut pour visibilite de bout en bout",
]
txBox = slide.shapes.add_textbox(Inches(7.2), Inches(4.9), Inches(5), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(rpa_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Segoe UI"
    p.space_after = Pt(3)


# ============================================================
# SLIDE 8 — TABLEAUX DE BORD & ANALYTIQUE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), RGBColor(0x6C, 0x5C, 0xE7))
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.6),
             "Tableaux de Bord & Analytique", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(1.1), Inches(1.05), Inches(10), Inches(0.4),
             '"Des problemes bruts aux insights executifs en un clic."',
             font_size=14, color=RGBColor(0x6C, 0x5C, 0xE7), bold=True)

dashboards = [
    ("Tableau de Bord Executif", [
        "Total clients par type",
        "Compteurs et tendances anomalies",
        "Vue d'ensemble FATCA",
        "Tickets en attente",
        "Metriques taux de correction",
        "Widget anomalies recentes"
    ], ACCENT_BLUE),
    ("Tableau de Bord KPI", [
        "Indicateurs cles de performance",
        "KPI specifiques par agence",
        "Analyse historique",
        "Filtrage par plage de dates",
        "Visualisation des tendances",
        "Indicateurs de statut"
    ], RGBColor(0x6C, 0x5C, 0xE7)),
    ("Performance Agences", [
        "Distribution anomalies par agence",
        "Stats de correction par agence",
        "Tendances hebdomadaires",
        "Meilleurs/moins bons performers",
        "Comparaison inter-agences",
        "Metriques de responsabilite"
    ], ACCENT_TEAL),
    ("Rapports & Exports", [
        "Rapports qualite hebdomadaires",
        "Generation PDF",
        "Export tableau de bord Excel",
        "Export donnees CSV",
        "Statistiques de synthese",
        "Reporting par date"
    ], GOLD),
]

for i, (title, items, color) in enumerate(dashboards):
    x = Inches(0.8 + i * 3.1)
    y = Inches(1.7)
    add_card(slide, x, y, Inches(2.8), Inches(4.0))
    add_accent_bar(slide, x, y, Inches(2.8), Inches(0.06), color)
    add_text_box(slide, x + Inches(0.25), y + Inches(0.2), Inches(2.3), Inches(0.4),
                 title, font_size=15, color=color, bold=True)
    txBox = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.7), Inches(2.3), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True
    for j, item in enumerate(items):
        p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
        p.text = f"  {item}"
        p.font.size = Pt(11)
        p.font.color.rgb = LIGHT_GRAY
        p.font.name = "Segoe UI"
        p.space_after = Pt(5)

metrics = [
    ("Temps Reel", "Donnees en direct", ACCENT_BLUE),
    ("Multi-Format", "PDF / Excel / CSV", RGBColor(0x6C, 0x5C, 0xE7)),
    ("Multi-Agences", "Vues par agence", ACCENT_TEAL),
    ("Historique", "Analyse de tendances", GOLD),
]
for i, (title, desc, color) in enumerate(metrics):
    x = Inches(0.8 + i * 3.1)
    y = Inches(6.0)
    add_card(slide, x, y, Inches(2.8), Inches(0.9))
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(2.4), Inches(0.35),
                 title, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.45), Inches(2.4), Inches(0.3),
                 desc, font_size=10, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 9 — DETECTION DOUBLONS & GESTION DONNEES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), RGBColor(0xE8, 0x4D, 0x8A))
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.6),
             "Detection de Doublons & Gestion des Donnees", font_size=30, color=WHITE, bold=True)

# Duplicates
add_card(slide, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2))
add_text_box(slide, Inches(1.1), Inches(1.7), Inches(5), Inches(0.4),
             '"Un client, un dossier — eliminez les doublons couteux."',
             font_size=14, color=RGBColor(0xE8, 0x4D, 0x8A), bold=True)

dup_features = [
    ("Detection Algorithmique", "Des algorithmes avances trouvent les doublons que les humains manquent dans toute la base client"),
    ("Score de Confiance", "Niveaux Haut / Moyen / Bas — priorisez les correspondances a haute confiance"),
    ("Confirmation Manuelle", "Aucune fusion erronee — validation humaine avant toute action de fusion"),
    ("Filtrage par Type", "Vues separees des doublons pour particuliers et entreprises"),
    ("Tableau de Bord Stats", "Suivi total, en attente, confirmes et progression de la deduplication"),
]

for i, (title, desc) in enumerate(dup_features):
    y = Inches(2.3 + i * 0.85)
    add_accent_bar(slide, Inches(1.1), y + Inches(0.05), Inches(0.06), Inches(0.5), RGBColor(0xE8, 0x4D, 0x8A))
    add_text_box(slide, Inches(1.3), y, Inches(5), Inches(0.3), title, font_size=13, color=WHITE, bold=True)
    add_text_box(slide, Inches(1.3), y + Inches(0.3), Inches(5), Inches(0.4), desc, font_size=10, color=LIGHT_GRAY)

# Data import
add_card(slide, Inches(6.9), Inches(1.5), Inches(5.6), Inches(2.5))
add_text_box(slide, Inches(7.2), Inches(1.7), Inches(5), Inches(0.4),
             "Import de Donnees Flexible", font_size=18, color=ACCENT_BLUE, bold=True)
import_items = [
    "Upload CSV avec validation automatique",
    "Traitement de fichiers Excel",
    "Historique d'upload avec suivi des erreurs",
    "Monitoring du temps d'execution",
    "Traitement par lots des enregistrements",
]
txBox = slide.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(import_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Segoe UI"
    p.space_after = Pt(3)

# Security
add_card(slide, Inches(6.9), Inches(4.2), Inches(5.6), Inches(2.5))
add_text_box(slide, Inches(7.2), Inches(4.4), Inches(5), Inches(0.4),
             "Securite Entreprise", font_size=18, color=ACCENT_TEAL, bold=True)
add_text_box(slide, Inches(7.2), Inches(4.7), Inches(5), Inches(0.3),
             '"Securite integree — pas ajoutee apres coup."', font_size=11, color=ACCENT_TEAL)
sec_items = [
    "Keycloak SSO (OAuth2) — connexion unique",
    "Synchronisation LDAP/Active Directory",
    "Cookies de session HttpOnly (anti-XSS)",
    "Journalisation complete (IP, user-agent, horodatage)",
    "Hachage mot de passe bcrypt",
    "90+ points API REST pour l'integration",
]
txBox = slide.shapes.add_textbox(Inches(7.2), Inches(5.1), Inches(5), Inches(1.5))
tf = txBox.text_frame
tf.word_wrap = True
for i, item in enumerate(sec_items):
    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
    p.text = f"  {item}"
    p.font.size = Pt(11)
    p.font.color.rgb = LIGHT_GRAY
    p.font.name = "Segoe UI"
    p.space_after = Pt(3)


# ============================================================
# SLIDE 10 — ARCHITECTURE MULTI-AGENCES
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), ACCENT_TEAL)
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.6),
             "Architecture Multi-Agences", font_size=30, color=WHITE, bold=True)
add_text_box(slide, Inches(1.1), Inches(1.05), Inches(10), Inches(0.4),
             '"Scalable de 1 agence a 500 — meme plateforme."', font_size=14, color=ACCENT_TEAL, bold=True)

# HQ
add_card(slide, Inches(3.5), Inches(1.7), Inches(6.3), Inches(1.6))
add_accent_bar(slide, Inches(3.5), Inches(1.7), Inches(6.3), Inches(0.06), ACCENT_BLUE)
add_text_box(slide, Inches(3.8), Inches(1.9), Inches(5.5), Inches(0.4),
             "SIEGE — Vue Globale", font_size=16, color=ACCENT_BLUE, bold=True)
add_text_box(slide, Inches(3.8), Inches(2.3), Inches(5.5), Inches(0.7),
             "Tableaux de bord inter-agences  |  KPI globaux  |  Toutes anomalies  |  Gestion utilisateurs\nSupervision FATCA  |  Configuration systeme  |  Generation de rapports",
             font_size=11, color=LIGHT_GRAY)

branches = [
    ("Agence A", "Donnees de l'agence\nSes anomalies\nCorrections locales\nFATCA local"),
    ("Agence B", "Donnees de l'agence\nSes anomalies\nCorrections locales\nFATCA local"),
    ("Agence C", "Donnees de l'agence\nSes anomalies\nCorrections locales\nFATCA local"),
    ("Agence N...", "Donnees de l'agence\nSes anomalies\nCorrections locales\nFATCA local"),
]

for i, (name, desc) in enumerate(branches):
    x = Inches(0.8 + i * 3.1)
    y = Inches(3.7)
    add_card(slide, x, y, Inches(2.8), Inches(2.4))
    add_accent_bar(slide, x, y, Inches(2.8), Inches(0.06), ACCENT_TEAL)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.2), Inches(2.2), Inches(0.35),
                 name, font_size=14, color=ACCENT_TEAL, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.6), Inches(2.2), Inches(1.5),
                 desc, font_size=10, color=LIGHT_GRAY)

for i in range(4):
    x = Inches(2.2 + i * 3.1)
    add_text_box(slide, x, Inches(3.3), Inches(1.0), Inches(0.4),
                 "|", font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_card(slide, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9))
benefits = ["Zero config par agence", "Isolation automatique des donnees", "Responsabilite par agence", "Visibilite totale siege"]
for i, b in enumerate(benefits):
    x = Inches(1.2 + i * 3.0)
    add_accent_bar(slide, x, Inches(6.55), Inches(0.06), Inches(0.4), ACCENT_TEAL)
    add_text_box(slide, x + Inches(0.2), Inches(6.5), Inches(2.5), Inches(0.5),
                 b, font_size=13, color=WHITE, bold=True)


# ============================================================
# SLIDE 11 — TOP 5 ARGUMENTS DE VENTE
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide, SECTION_BG)

add_accent_bar(slide, Inches(0.8), Inches(0.5), Inches(0.08), Inches(0.5), GOLD)
add_text_box(slide, Inches(1.1), Inches(0.5), Inches(10), Inches(0.6),
             "Top 5 Raisons de Choisir BanK Data Quality Monitor", font_size=30, color=WHITE, bold=True)

selling_points = [
    ("Conformite Reglementaire", "Filtrage FATCA + validation 4-yeux + piste d'audit complete — toujours pret pour l'inspection. Respectez les exigences IRS et reglementaires locales sans effort manuel.", GOLD),
    ("Automatisation de Bout en Bout", "De la detection a la reconciliation — pipeline entierement automatise. L'integration RPA elimine la re-saisie manuelle entre les systemes.", ACCENT_BLUE),
    ("Integration Core Banking Directe", "Connexion directe a Informix/BSIC via JDBC — pas de middleware, zero silo de donnees. Les corrections alimentent directement le systeme source.", ORANGE),
    ("Scalabilite Multi-Agences", "Responsabilite au niveau agence avec visibilite du siege. Meme plateforme pour 1 ou 500+ agences, zero configuration par agence.", ACCENT_TEAL),
    ("ROI Mesurable", "Tableaux de bord KPI qui prouvent l'amelioration de la qualite des donnees. Suivez les taux de correction, la performance des agences et les metriques de conformite.", RGBColor(0x6C, 0x5C, 0xE7)),
]

for i, (title, desc, color) in enumerate(selling_points):
    y = Inches(1.4 + i * 1.15)

    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.05), Inches(0.6), Inches(0.6))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = "Segoe UI"
    p.alignment = PP_ALIGN.CENTER

    add_text_box(slide, Inches(1.7), y, Inches(10), Inches(0.4),
                 title, font_size=18, color=color, bold=True)
    add_text_box(slide, Inches(1.7), y + Inches(0.4), Inches(10), Inches(0.5),
                 desc, font_size=12, color=LIGHT_GRAY)


# ============================================================
# SLIDE 12 — MERCI / CONTACT
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_dark_bg(slide, SECTION_BG)

add_accent_bar(slide, Inches(4.8), Inches(2.2), Inches(3.7), Inches(0.06), ACCENT_TEAL)

add_text_box(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0),
             "Pret a Transformer Votre Qualite de Donnees ?",
             font_size=40, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(0.8),
             "Laissez-nous vous montrer comment BanK Data Quality Monitor peut aider votre\netablissement a atteindre la conformite reglementaire et l'excellence operationnelle.",
             font_size=16, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_card(slide, Inches(4.0), Inches(4.8), Inches(5.3), Inches(1.5))
add_text_box(slide, Inches(4.3), Inches(5.0), Inches(4.7), Inches(0.4),
             "Demander une Demo", font_size=20, color=ACCENT_TEAL, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(4.3), Inches(5.5), Inches(4.7), Inches(0.6),
             "Contactez notre equipe pour une demonstration personnalisee\nadaptee a votre environnement bancaire.",
             font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(2), Inches(6.6), Inches(9.3), Inches(0.5),
             "BanK Data Quality Monitor — Detecter. Corriger. Se Conformer.",
             font_size=14, color=ACCENT_BLUE, bold=True, alignment=PP_ALIGN.CENTER)


# === SAVE ===
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "BanK_DataQuality_Marketing_Presentation.pptx")
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
